import os
import tempfile
import uuid
import subprocess
from pathlib import Path
from fastapi import FastAPI, UploadFile, File, Form, Request
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from pydantic import BaseModel
from transformers import pipeline
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog, simpledialog, messagebox
import comtypes.client  # MS PowerPoint Live Detection
# ppt.py

from fastapi import APIRouter

router = APIRouter()

@router.get("/ppt-example")
def test_ppt():
    return {"message": "PPT route is working!"}

app = FastAPI(title="📽️ Smart PPT AI Assistant")
TEMP_DIR = Path(tempfile.gettempdir())
SUPPORTED_FORMATS = ["pptx"]

qa_model = pipeline("text2text-generation", model="google/flan-t5-small")

# 🧠 Prompt-based content processing
def generate_content(topic: str, style="structured") -> str:
    prompt = f"Create a detailed and {style} slide presentation on: {topic}"
    return qa_model(prompt, max_length=1024)[0]["generated_text"]

# 📊 Create PPT from content
def create_ppt(content: str, output_path: str) -> str:
    ppt = Presentation()
    for section in content.strip().split("\n\n"):
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        lines = section.strip().split("\n")
        if lines:
            slide.shapes.title.text = lines[0]
            if len(lines) > 1:
                slide.placeholders[1].text = "\n".join(lines[1:])
    ppt.save(output_path)
    return output_path

# 🧾 Detect open PowerPoint file
def get_open_pptx_path() -> str:
    try:
        app = comtypes.client.CreateObject("PowerPoint.Application")
        pres = app.Presentations
        if pres.Count > 0:
            return pres.Item(1).FullName
    except Exception:
        pass
    return None

# 🌐 FastAPI Routes
@app.get("/", response_class=HTMLResponse)
async def main_page():
    return HTMLResponse("""
    <html><body>
    <h2>📽️ AI Slide Assistant</h2>
    <form action="/interactive/" method="post">
        <input type="text" name="topic" placeholder="Enter your topic" style="width: 400px;"><br><br>
        <input type="submit" value="Generate PPT">
    </form></body></html>
    """)

@app.post("/interactive/")
async def interactive_generate(request: Request):
    form = await request.form()
    topic = form.get("topic")
    result = generate_content(topic)
    output_path = TEMP_DIR / f"{uuid.uuid4()}.pptx"
    create_ppt(result, output_path)
    return FileResponse(output_path, filename="interactive_slides.pptx")

@app.post("/upload/")
async def upload_ppt(file: UploadFile = File(...), task: str = Form("Summarize into slides")):
    ext = file.filename.split(".")[-1].lower()
    if ext not in SUPPORTED_FORMATS:
        return JSONResponse(content={"error": "Unsupported format"}, status_code=400)

    file_path = TEMP_DIR / f"{uuid.uuid4()}.pptx"
    with open(file_path, "wb") as f:
        f.write(await file.read())

    prs = Presentation(file_path)
    content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    result = qa_model(f"{task}:\n{content}", max_length=512)[0]["generated_text"]
    output_path = TEMP_DIR / f"processed_{uuid.uuid4()}.pptx"
    create_ppt(result, output_path)
    return FileResponse(output_path, filename="processed_slides.pptx")

@app.post("/process-open-ppt/")
async def process_open_ppt(task: str = Form("Summarize this presentation")):
    path = get_open_pptx_path()
    if not path or not os.path.exists(path):
        return JSONResponse(content={"error": "No open PowerPoint file detected"}, status_code=404)

    prs = Presentation(path)
    content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    result = qa_model(f"{task}:\n{content}", max_length=512)[0]["generated_text"]
    output_path = TEMP_DIR / f"auto_processed_{uuid.uuid4()}.pptx"
    create_ppt(result, output_path)
    return FileResponse(output_path, filename=output_path.name)

@app.post("/generate/")
async def generate_from_topic(topic: str = Form(...)):
    content = generate_content(topic)
    output_path = TEMP_DIR / f"generated_{uuid.uuid4()}.pptx"
    create_ppt(content, output_path)
    return FileResponse(output_path, filename="generated_output.pptx")

# 🖥️ CLI Mode
def cli_runner():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--file", help="Path to PPTX file")
    parser.add_argument("--task", help="Instruction to process slides")
    parser.add_argument("--topic", help="Topic to generate new PPT")
    args = parser.parse_args()

    if args.topic:
        content = generate_content(args.topic)
        output = create_ppt(content, "cli_generated_output.pptx")
        print(f"✅ Created: {output}")
        return

    file = args.file or get_open_pptx_path()
    if file and file.endswith(".pptx"):
        prs = Presentation(file)
        content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        result = qa_model(f"{args.task}:\n{content}", max_length=512)[0]["generated_text"]
        output = create_ppt(result, "cli_processed_output.pptx")
        print(f"✅ Done: {output}")
    else:
        print("❌ No valid .pptx file found.")

# 🖱️ GUI Mode
def start_gui():
    root = tk.Tk()
    root.title("Smart PPT Assistant")

    def run():
        topic = simpledialog.askstring("Topic", "Enter presentation topic:")
        if topic:
            content = generate_content(topic)
            path = create_ppt(content, f"gui_ppt_{uuid.uuid4()}.pptx")
            messagebox.showinfo("Done", f"Saved: {path}")

    tk.Button(root, text="Generate Slides", command=run, height=2, width=30).pack(pady=20)
    root.mainloop()

if __name__ == "__main__":
    import sys
    if "--gui" in sys.argv:
        start_gui()
    elif len(sys.argv) > 1:
        cli_runner()
    else:
        import uvicorn
        uvicorn.run(app, host="0.0.0.0", port=8000)
